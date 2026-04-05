"""PPTX（Googleスライド互換）レポート生成モジュール

analyze_periodの結果とAIコメンタリーからPowerPointプレゼンテーションを生成する。
Googleスライドで直接開ける.pptx形式で出力する。
"""

from __future__ import annotations

import logging
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ── カラーパレット（洗練されたモノトーン + アクセント） ──
BG_DARK = RGBColor(0x11, 0x14, 0x1A)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)
TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
SURFACE = RGBColor(0xF9, 0xFA, 0xFB)
ACCENT = RGBColor(0x37, 0x4B, 0xF5)   # インディゴブルー
ACCENT_LIGHT = RGBColor(0xEE, 0xF0, 0xFF)
POSITIVE = RGBColor(0x05, 0x96, 0x69)
NEGATIVE = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM = RGBColor(0xF5, 0x9E, 0x0B)     # アンバー
CORAL = RGBColor(0xEF, 0x44, 0x44)    # コーラル

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _n(v) -> str:
    """数値フォーマット"""
    if v is None:
        return "--"
    if isinstance(v, float) and v == int(v):
        v = int(v)
    if isinstance(v, int):
        if v >= 10000:
            return f"{v / 10000:.1f}万"
        return f"{v:,}"
    return str(v)


def _pct(v) -> str:
    if v is None:
        return "--"
    return f"{'+' if v > 0 else ''}{v:.1f}%"


def _set_cell(cell, text: str, size: int = 10, bold: bool = False,
              color: RGBColor = TEXT_PRIMARY, align=PP_ALIGN.LEFT,
              bg: RGBColor | None = None):
    """テーブルセルの設定"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(1)
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Emu(63500)
    cell.margin_right = Emu(63500)
    cell.margin_top = Emu(36576)
    cell.margin_bottom = Emu(36576)
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _slide_bg(slide, color: RGBColor = SURFACE):
    """スライド背景色"""
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def _page_header(slide, cn: str, period: str):
    """ページヘッダー（ミニマル）"""
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(3), Inches(0.3))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = cn
    r.font.size = Pt(9)
    r.font.color.rgb = TEXT_MUTED
    r.font.name = "Helvetica Neue"
    r.font.bold = True

    tb2 = slide.shapes.add_textbox(Inches(9.5), Inches(0.3), Inches(3.2), Inches(0.3))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run()
    r2.text = period
    r2.font.size = Pt(9)
    r2.font.color.rgb = TEXT_MUTED
    r2.font.name = "Helvetica Neue"


def _section_title(slide, text: str, top: float = 0.85):
    """セクションタイトル（左アクセントドット付き）"""
    # アクセントドット
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.8), Inches(top + 0.12), Pt(8), Pt(8),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1.1), Inches(top), Inches(10), Inches(0.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TEXT_PRIMARY
    r.font.name = "Helvetica Neue"


def _kpi_card(slide, left: float, top: float, w: float, h: float,
              label: str, value: str, sub: str = "", accent: RGBColor = ACCENT):
    """KPIカード（角丸白カード）"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)
    # 左上にアクセントライン
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Pt(3), Inches(h),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # ラベル
    lb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.12), Inches(w - 0.3), Inches(0.2))
    lr = lb.text_frame.paragraphs[0].add_run()
    lr.text = label
    lr.font.size = Pt(8)
    lr.font.color.rgb = TEXT_MUTED
    lr.font.name = "Helvetica Neue"
    lr.font.bold = True

    # 値
    vb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.35), Inches(w - 0.3), Inches(0.35))
    vr = vb.text_frame.paragraphs[0].add_run()
    vr.text = value
    vr.font.size = Pt(20)
    vr.font.bold = True
    vr.font.color.rgb = TEXT_PRIMARY
    vr.font.name = "Helvetica Neue"

    # サブテキスト
    if sub:
        sb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.7), Inches(w - 0.3), Inches(0.2))
        sr = sb.text_frame.paragraphs[0].add_run()
        sr.text = f"前月比 {sub}"
        sr.font.size = Pt(8)
        sr.font.bold = True
        sr.font.color.rgb = POSITIVE if sub.startswith("+") else NEGATIVE if sub.startswith("-") else TEXT_MUTED
        sr.font.name = "Helvetica Neue"


def _ai_text_slide(prs, blank, cn: str, period: str,
                   title: str, body: str, accent: RGBColor = ACCENT):
    """AI考察テキストスライド"""
    slide = prs.slides.add_slide(blank)
    _slide_bg(slide, SURFACE)
    _page_header(slide, cn, period)
    _section_title(slide, title)

    # メインテキストカード
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # 上部アクセントライン
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(11.7), Pt(3),
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = accent
    top_line.line.fill.background()

    # テキスト本文
    tb = slide.shapes.add_textbox(
        Inches(1.3), Inches(1.85), Inches(10.7), Inches(4.9),
    )
    tf = tb.text_frame
    tf.word_wrap = True

    # 長いテキストは段落分割
    paragraphs = body.split("\n") if body else [""]
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.line_spacing = Pt(18)
        run = p.add_run()
        run.text = para_text.strip()
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_PRIMARY
        run.font.name = "Helvetica Neue"

    return slide


def generate_pptx(context: dict, output_path: Path) -> Path:
    """レポートのPPTXを生成する"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    cn = context.get("client_name", "")
    period = context.get("period", "")
    ai = context.get("ai_commentary", {})

    # ================================================================
    # P1: 表紙
    # ================================================================
    s1 = prs.slides.add_slide(blank)
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK

    # 左アクセントライン
    bar = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.8), Inches(0.8), Pt(3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # サブタイトル
    st = s1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8), Inches(0.4))
    sr = st.text_frame.paragraphs[0].add_run()
    sr.text = "TikTok Monthly Report"
    sr.font.size = Pt(14)
    sr.font.color.rgb = TEXT_MUTED
    sr.font.name = "Helvetica Neue"
    sr.font.letter_spacing = Pt(2)

    # クライアント名
    t1 = s1.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(10), Inches(1.0))
    tr1 = t1.text_frame.paragraphs[0].add_run()
    tr1.text = cn
    tr1.font.size = Pt(40)
    tr1.font.color.rgb = WHITE
    tr1.font.bold = True
    tr1.font.name = "Helvetica Neue"

    # 期間
    p1 = s1.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8), Inches(0.5))
    pr1 = p1.text_frame.paragraphs[0].add_run()
    pr1.text = period
    pr1.font.size = Pt(16)
    pr1.font.color.rgb = TEXT_MUTED
    pr1.font.name = "Helvetica Neue"

    # 月間総括
    overall = ai.get("overall_assessment", "")
    if overall:
        # 仕切り線
        div = s1.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(4.7), Inches(4), Pt(1),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = RGBColor(0x30, 0x35, 0x45)
        div.line.fill.background()

        o1 = s1.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(10), Inches(1.2))
        otf = o1.text_frame
        otf.word_wrap = True
        orun = otf.paragraphs[0].add_run()
        orun.text = overall
        orun.font.size = Pt(12)
        orun.font.color.rgb = RGBColor(0x8B, 0x92, 0xA5)
        orun.font.name = "Helvetica Neue"

    # 右下ブランド
    br = s1.shapes.add_textbox(Inches(10.5), Inches(6.6), Inches(2.3), Inches(0.4))
    brf = br.text_frame
    brf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    brr = brf.paragraphs[0].add_run()
    brr.text = "LEAD ONE"
    brr.font.size = Pt(10)
    brr.font.color.rgb = RGBColor(0x40, 0x45, 0x55)
    brr.font.bold = True
    brr.font.name = "Helvetica Neue"
    brr.font.letter_spacing = Pt(3)

    # ================================================================
    # P2: KPIダッシュボード
    # ================================================================
    s2 = prs.slides.add_slide(blank)
    _slide_bg(s2, SURFACE)
    _page_header(s2, cn, period)
    _section_title(s2, "KPIダッシュボード")

    mom_views = context.get("mom_views")
    mom_likes = context.get("mom_likes")
    mom_comments = context.get("mom_comments")
    mom_shares = context.get("mom_shares")

    # メインKPIカード（上段4枚）
    main_kpis = [
        ("総再生数", _n(context.get("total_views")),
         _pct(mom_views) if mom_views is not None else "", ACCENT),
        ("いいね", _n(context.get("total_likes")),
         _pct(mom_likes) if mom_likes is not None else "", CORAL),
        ("コメント", _n(context.get("total_comments")),
         _pct(mom_comments) if mom_comments is not None else "", WARM),
        ("シェア", _n(context.get("total_shares")),
         _pct(mom_shares) if mom_shares is not None else "", POSITIVE),
    ]
    for i, (label, value, sub, color) in enumerate(main_kpis):
        _kpi_card(s2, 0.8 + i * 3.05, 1.5, 2.75, 0.95, label, value, sub, color)

    # サブKPIカード（中段）
    eng_rate = context.get("engagement_rate", 0)
    prof_rate = context.get("profile_transition_rate")
    fg = context.get("follower_growth")

    sub_kpis = [
        ("投稿数", f"{context.get('post_count', 0)}本", "", TEXT_PRIMARY),
        ("平均再生数/投稿", _n(context.get("avg_views_per_post")), "", TEXT_PRIMARY),
        ("エンゲージメント率", f"{eng_rate:.2f}%", "", ACCENT),
    ]
    if prof_rate is not None:
        sub_kpis.append(("プロフィール遷移率", f"{prof_rate:.2f}%", "", ACCENT))
    if fg:
        change = fg.get("change", 0)
        sub_kpis.append((
            "フォロワー増減",
            f"{'+'if change > 0 else ''}{_n(change)}",
            f"{fg['start_count']:,} → {fg['end_count']:,}",
            POSITIVE if change > 0 else NEGATIVE,
        ))

    sub_w = min(2.75, 12.1 / max(len(sub_kpis), 1) - 0.15)
    for i, (label, value, sub, color) in enumerate(sub_kpis):
        _kpi_card(s2, 0.8 + i * (sub_w + 0.15), 2.7, sub_w, 0.95, label, value, sub, color)

    # 詳細テーブル
    detail_label = s2.shapes.add_textbox(Inches(0.8), Inches(3.95), Inches(3), Inches(0.25))
    dlr = detail_label.text_frame.paragraphs[0].add_run()
    dlr.text = "DETAIL"
    dlr.font.size = Pt(8)
    dlr.font.color.rgb = TEXT_MUTED
    dlr.font.bold = True
    dlr.font.name = "Helvetica Neue"
    dlr.font.letter_spacing = Pt(2)

    rows_data = [
        ("指標", "今月実績", "前月比"),
        ("総再生数", _n(context.get("total_views")), _pct(mom_views)),
        ("いいね", _n(context.get("total_likes")), _pct(mom_likes)),
        ("コメント", _n(context.get("total_comments")), _pct(mom_comments)),
        ("シェア", _n(context.get("total_shares")), _pct(mom_shares)),
        ("プロフィール閲覧", _n(context.get("total_profile_views")), "--"),
        ("エンゲージメント率", f"{eng_rate:.2f}%", "--"),
        ("投稿数", str(context.get("post_count", 0)), "--"),
        ("平均再生数/投稿", _n(context.get("avg_views_per_post")), "--"),
    ]

    tbl = s2.shapes.add_table(
        len(rows_data), 3,
        Inches(0.8), Inches(4.25), Inches(11.7), Inches(3.0),
    ).table
    tbl.columns[0].width = Inches(4.5)
    tbl.columns[1].width = Inches(3.6)
    tbl.columns[2].width = Inches(3.6)

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            is_h = ri == 0
            c = WHITE if is_h else TEXT_PRIMARY
            bg_c = TEXT_PRIMARY if is_h else (SURFACE if ri % 2 == 0 else None)

            if ci == 2 and not is_h and val not in ("--", None):
                try:
                    num = float(val.replace("+", "").replace("%", ""))
                    c = POSITIVE if num > 0 else NEGATIVE if num < 0 else TEXT_MUTED
                except (ValueError, AttributeError):
                    pass

            _set_cell(tbl.cell(ri, ci), val or "--", size=10, bold=is_h or ci == 1,
                      color=c, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
                      bg=bg_c)

    # ================================================================
    # P3: 動画毎レポート（ページネーション）
    # ================================================================
    posts = context.get("posts", [])
    if posts:
        chunk_size = 12
        for page_start in range(0, len(posts), chunk_size):
            chunk = posts[page_start:page_start + chunk_size]
            s3 = prs.slides.add_slide(blank)
            _slide_bg(s3, SURFACE)
            _page_header(s3, cn, period)

            total_count = len(posts)
            if total_count > chunk_size:
                suffix = f" ({page_start + 1}-{page_start + len(chunk)} / {total_count})"
            else:
                suffix = f" ({total_count})"
            _section_title(s3, f"動画パフォーマンス{suffix}")

            cols = ["#", "タイトル", "投稿日", "再生数", "いいね",
                    "コメント", "シェア", "ENG率", "完了率", "2秒率"]
            col_w = [Inches(0.4), Inches(3.8), Inches(1.1), Inches(1.2),
                     Inches(1.0), Inches(1.0), Inches(1.0), Inches(0.9),
                     Inches(0.8), Inches(0.8)]

            tbl = s3.shapes.add_table(
                len(chunk) + 1, len(cols),
                Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5),
            ).table

            for ci, (header, w) in enumerate(zip(cols, col_w)):
                tbl.columns[ci].width = w
                _set_cell(tbl.cell(0, ci), header, size=9, bold=True,
                          color=WHITE, align=PP_ALIGN.CENTER if ci > 1 else PP_ALIGN.LEFT,
                          bg=TEXT_PRIMARY)

            for ri, post in enumerate(chunk, start=1):
                p = post if isinstance(post, SimpleNamespace) else SimpleNamespace(**post)
                rank = page_start + ri
                is_top = rank <= 3
                row_bg = ACCENT_LIGHT if is_top else (SURFACE if ri % 2 == 0 else None)

                caption = p.caption if len(p.caption) <= 30 else p.caption[:29] + "..."

                _set_cell(tbl.cell(ri, 0), str(rank), size=9, bold=is_top,
                          color=ACCENT if is_top else TEXT_MUTED, align=PP_ALIGN.CENTER, bg=row_bg)
                _set_cell(tbl.cell(ri, 1), caption, size=9, bold=is_top,
                          color=TEXT_PRIMARY, bg=row_bg)
                _set_cell(tbl.cell(ri, 2), getattr(p, "post_date_display", ""),
                          size=9, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)
                _set_cell(tbl.cell(ri, 3), _n(p.views), size=9, bold=True,
                          color=TEXT_PRIMARY, align=PP_ALIGN.RIGHT, bg=row_bg)
                _set_cell(tbl.cell(ri, 4), _n(p.likes), size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
                _set_cell(tbl.cell(ri, 5), _n(p.comments), size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
                _set_cell(tbl.cell(ri, 6), _n(p.shares), size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
                eng = getattr(p, "engagement_rate", 0) or 0
                _set_cell(tbl.cell(ri, 7), f"{eng:.1f}%", size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)
                wtr = getattr(p, "watch_through_rate", None)
                _set_cell(tbl.cell(ri, 8), f"{wtr}%" if wtr else "--", size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)
                svr = getattr(p, "two_sec_view_rate", None)
                _set_cell(tbl.cell(ri, 9), f"{svr}%" if svr else "--", size=9,
                          color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)

    # ================================================================
    # P4: トップ投稿ハイライト
    # ================================================================
    top_posts = context.get("top_posts", [])
    if top_posts:
        s4 = prs.slides.add_slide(blank)
        _slide_bg(s4, SURFACE)
        _page_header(s4, cn, period)
        _section_title(s4, f"Top {min(len(top_posts), 5)} ハイライト")

        for i, post in enumerate(top_posts[:5]):
            p = post if isinstance(post, SimpleNamespace) else SimpleNamespace(**post)
            col = i % 3
            row = i // 3
            left = 0.8 + col * 4.05
            top = 1.5 + row * 2.85

            # カード
            card = s4.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(3.75), Inches(2.55),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = BORDER
            card.line.width = Pt(0.5)

            # ランク + 再生数（ヘッダー行）
            rank_bg = s4.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(3.75), Inches(0.5),
            )
            rank_bg.fill.solid()
            rank_bg.fill.fore_color.rgb = TEXT_PRIMARY
            rank_bg.line.fill.background()

            rb = s4.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.08), Inches(0.8), Inches(0.35))
            rr = rb.text_frame.paragraphs[0].add_run()
            rr.text = f"#{i + 1}"
            rr.font.size = Pt(14)
            rr.font.color.rgb = ACCENT
            rr.font.bold = True
            rr.font.name = "Helvetica Neue"

            vb = s4.shapes.add_textbox(Inches(left + 1.0), Inches(top + 0.08), Inches(2.5), Inches(0.35))
            vtf = vb.text_frame
            vtf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            vr = vtf.paragraphs[0].add_run()
            vr.text = f"{_n(p.views)} views"
            vr.font.size = Pt(13)
            vr.font.bold = True
            vr.font.color.rgb = WHITE
            vr.font.name = "Helvetica Neue"

            # タイトル
            caption = p.caption if len(p.caption) <= 45 else p.caption[:44] + "..."
            tb = s4.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.65), Inches(3.3), Inches(0.7))
            ttf = tb.text_frame
            ttf.word_wrap = True
            tr = ttf.paragraphs[0].add_run()
            tr.text = caption
            tr.font.size = Pt(10)
            tr.font.color.rgb = TEXT_PRIMARY
            tr.font.name = "Helvetica Neue"

            # メトリクス
            metrics = [
                (f"{_n(p.likes)}", "いいね", CORAL),
                (f"{_n(p.comments)}", "コメント", WARM),
                (f"{_n(p.shares)}", "シェア", POSITIVE),
            ]
            for mi, (mval, mlabel, mclr) in enumerate(metrics):
                mx = left + 0.2 + mi * 1.2
                my = top + 1.5

                mv = s4.shapes.add_textbox(Inches(mx), Inches(my), Inches(1.1), Inches(0.25))
                mvr = mv.text_frame.paragraphs[0].add_run()
                mvr.text = mval
                mvr.font.size = Pt(12)
                mvr.font.bold = True
                mvr.font.color.rgb = mclr
                mvr.font.name = "Helvetica Neue"

                ml = s4.shapes.add_textbox(Inches(mx), Inches(my + 0.22), Inches(1.1), Inches(0.18))
                mlr = ml.text_frame.paragraphs[0].add_run()
                mlr.text = mlabel
                mlr.font.size = Pt(7)
                mlr.font.color.rgb = TEXT_MUTED
                mlr.font.name = "Helvetica Neue"

            # ENG率 + 日付
            eng = getattr(p, "engagement_rate", 0) or 0
            info = s4.shapes.add_textbox(Inches(left + 0.2), Inches(top + 2.1), Inches(3.3), Inches(0.3))
            inf = info.text_frame.paragraphs[0].add_run()
            inf.text = f"ENG {eng:.1f}%  |  {getattr(p, 'post_date_display', '')}"
            inf.font.size = Pt(8)
            inf.font.color.rgb = TEXT_MUTED
            inf.font.name = "Helvetica Neue"

    # ================================================================
    # P4b: ワースト投稿分析
    # ================================================================
    worst_posts = context.get("worst_posts", [])
    if worst_posts:
        s4b = prs.slides.add_slide(blank)
        _slide_bg(s4b, SURFACE)
        _page_header(s4b, cn, period)
        _section_title(s4b, "改善対象の投稿")

        sub = s4b.shapes.add_textbox(Inches(1.1), Inches(1.3), Inches(8), Inches(0.25))
        sr = sub.text_frame.paragraphs[0].add_run()
        sr.text = "再生数下位 - 改善のヒントを見つける"
        sr.font.size = Pt(9)
        sr.font.color.rgb = TEXT_MUTED
        sr.font.name = "Helvetica Neue"

        cols = ["タイトル", "投稿日", "再生数", "いいね", "コメント", "シェア", "ENG率"]
        col_w = [Inches(4.5), Inches(1.2), Inches(1.3), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.1)]
        tbl = s4b.shapes.add_table(
            len(worst_posts) + 1, len(cols),
            Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.8),
        ).table

        for ci, (h, w) in enumerate(zip(cols, col_w)):
            tbl.columns[ci].width = w
            _set_cell(tbl.cell(0, ci), h, size=9, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT, bg=TEXT_PRIMARY)

        for ri, post in enumerate(worst_posts, start=1):
            p = post if isinstance(post, SimpleNamespace) else SimpleNamespace(**post)
            row_bg = SURFACE if ri % 2 == 0 else None
            _set_cell(tbl.cell(ri, 0), p.caption[:35] + ("..." if len(p.caption) > 35 else ""),
                      size=9, color=TEXT_PRIMARY, bg=row_bg)
            _set_cell(tbl.cell(ri, 1), getattr(p, "post_date_display", ""),
                      size=9, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)
            _set_cell(tbl.cell(ri, 2), _n(p.views), size=9, bold=True,
                      color=NEGATIVE, align=PP_ALIGN.RIGHT, bg=row_bg)
            _set_cell(tbl.cell(ri, 3), _n(p.likes), size=9,
                      color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
            _set_cell(tbl.cell(ri, 4), _n(p.comments), size=9,
                      color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
            _set_cell(tbl.cell(ri, 5), _n(p.shares), size=9,
                      color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, bg=row_bg)
            eng = getattr(p, "engagement_rate", 0) or 0
            _set_cell(tbl.cell(ri, 6), f"{eng:.1f}%", size=9,
                      color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, bg=row_bg)

    # ================================================================
    # P5: 数値報告（月別推移）
    # ================================================================
    monthly = context.get("monthly_transition", [])
    if monthly:
        s_num = prs.slides.add_slide(blank)
        _slide_bg(s_num, SURFACE)
        _page_header(s_num, cn, period)
        _section_title(s_num, "数値報告")

        # 表示する月数（最大12ヶ月目まで）
        months = monthly[:12]
        n_months = len(months)
        n_cols = n_months + 2  # カテゴリ列 + 指標列 + N月分

        # 行定義: カテゴリ / 指標名
        row_defs = [
            ("当月", "フォロワー数"),
            ("累計", "再生回数"),
            ("累計", "プロフィール表示回数"),
            ("累計", "フォロワー増加数"),
            ("月間", "再生回数"),
            ("月間", "フォロワー増加数"),
            ("月間", "プロフィール表示回数"),
            ("月間", "プロフィール遷移率"),
        ]

        n_rows = len(row_defs) + 1  # +1 for header
        tbl_w = 11.7
        tbl = s_num.shapes.add_table(
            n_rows, n_cols,
            Inches(0.8), Inches(1.5), Inches(tbl_w), Inches(5.5),
        ).table

        # 列幅設定
        cat_w = Inches(0.9)
        metric_w = Inches(1.5)
        month_w = Inches((tbl_w - 0.9 - 1.5) / max(n_months, 1))
        tbl.columns[0].width = cat_w
        tbl.columns[1].width = metric_w
        for ci in range(2, n_cols):
            tbl.columns[ci].width = month_w

        # ヘッダー行
        _set_cell(tbl.cell(0, 0), "", size=9, bold=True, color=WHITE, bg=TEXT_PRIMARY)
        _set_cell(tbl.cell(0, 1), "運用月", size=9, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, bg=TEXT_PRIMARY)
        for ci, m in enumerate(months):
            _set_cell(tbl.cell(0, ci + 2), f"{m['month_num']}ヶ月目", size=9, bold=True,
                      color=WHITE, align=PP_ALIGN.CENTER, bg=TEXT_PRIMARY)

        # カテゴリセル結合用: 同じカテゴリが連続する行を検出
        prev_cat = None
        cat_start = 1
        cat_groups: list[tuple[str, int, int]] = []
        for ri, (cat, _) in enumerate(row_defs, start=1):
            if cat != prev_cat:
                if prev_cat is not None:
                    cat_groups.append((prev_cat, cat_start, ri - 1))
                cat_start = ri
                prev_cat = cat
        if prev_cat is not None:
            cat_groups.append((prev_cat, cat_start, len(row_defs)))

        # カテゴリ列 + 指標列を設定
        for cat, start_r, end_r in cat_groups:
            # カテゴリセルに値を設定（最初の行のみ、残りは空）
            for ri in range(start_r, end_r + 1):
                cat_bg = RGBColor(0xF0, 0xF1, 0xF3)
                if ri == start_r:
                    _set_cell(tbl.cell(ri, 0), cat, size=9, bold=True,
                              color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, bg=cat_bg)
                else:
                    _set_cell(tbl.cell(ri, 0), "", size=9, bg=cat_bg)

        # 指標名列
        for ri, (_, metric) in enumerate(row_defs, start=1):
            _set_cell(tbl.cell(ri, 1), metric, size=9, bold=True,
                      color=TEXT_PRIMARY, align=PP_ALIGN.CENTER,
                      bg=RGBColor(0xF0, 0xF1, 0xF3))

        # データ列
        cum_views = 0
        for ci, m in enumerate(months):
            col = ci + 2
            mv = m.get("monthly_views", 0)
            cum_views_val = m.get("cumulative_views", 0)
            row_bg = SURFACE if ci % 2 == 0 else None

            data_vals = [
                "--",                          # 当月フォロワー数（要手入力）
                _n(cum_views_val),             # 累計再生回数
                "--",                          # 累計プロフィール表示回数
                "--",                          # 累計フォロワー増加数
                _n(mv),                        # 月間再生回数
                "--",                          # 月間フォロワー増加数
                "--",                          # 月間プロフィール表示回数
                "--",                          # 月間プロフィール遷移率
            ]

            for ri, val in enumerate(data_vals):
                _set_cell(tbl.cell(ri + 1, col), val, size=9,
                          color=TEXT_PRIMARY if val != "--" else TEXT_MUTED,
                          align=PP_ALIGN.CENTER, bg=row_bg)

        # カテゴリセル結合（python-pptxでは垂直結合をサポート）
        for cat, start_r, end_r in cat_groups:
            if end_r > start_r:
                tbl.cell(start_r, 0).merge(tbl.cell(end_r, 0))

    # ================================================================
    # P6: 総評・パフォーマンス分析（AI）
    # ================================================================
    _ai_text_slide(prs, blank, cn, period,
                   "総評・パフォーマンス分析",
                   ai.get("best_post_analysis", ""),
                   accent=ACCENT)

    # ================================================================
    # P7: 改善提案（AI）
    # ================================================================
    _ai_text_slide(prs, blank, cn, period,
                   "改善提案",
                   ai.get("improvement_suggestions", ""),
                   accent=WARM)

    # ================================================================
    # P8: 来月のアクションプラン（AI）
    # ================================================================
    _ai_text_slide(prs, blank, cn, period,
                   "来月のアクションプラン",
                   ai.get("next_month_plan", ""),
                   accent=POSITIVE)

    # ================================================================
    # 保存
    # ================================================================
    prs.save(str(output_path))
    logger.info("PPTX生成完了: %s", output_path)
    return output_path
